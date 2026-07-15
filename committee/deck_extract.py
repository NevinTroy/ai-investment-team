"""Text extraction from uploaded pitch decks (.pptx / .pdf).

Turns raw uploaded bytes into a structured, LLM- and UI-readable payload:

    {
      "kind": "pptx" | "pdf",
      "filename": "seed-deck.pptx",
      "sections": [{"label": "Slide 1", "text": "..."}, ...],
      "text": "<all sections joined>",
      "char_count": 1234,
    }

``sections`` preserves slide/page boundaries so the frontend can show the
extracted content the way it was laid out. Extraction is best-effort per
slide/page — a single unreadable slide is skipped rather than failing the whole
deck — but a file that yields no text at all (or an unsupported type) raises
``ValueError`` so the caller can surface a clean message.
"""

from __future__ import annotations

import io
import logging
import re
from pathlib import Path

logger = logging.getLogger("committee.deck_extract")

SUPPORTED_EXTS = {".pptx", ".pdf"}

# Guardrail for the identification LLM call — a long deck's full text can blow
# the prompt budget, so callers truncate to this before sending to the model.
MAX_LLM_CHARS = 24000


def _clean(text: str) -> str:
    """Collapse runs of blank lines/trailing spaces so extracted text stays compact."""
    lines = [ln.rstrip() for ln in (text or "").splitlines()]
    out: list[str] = []
    for ln in lines:
        if ln.strip() or (out and out[-1].strip()):
            out.append(ln)
    return "\n".join(out).strip()


def _extract_pptx(data: bytes) -> tuple[list[dict], int]:
    """Returns (sections-with-text, total slide count)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)
    sections: list[dict] = []
    for i, slide in enumerate(slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            parts.append(" | ".join(cells))
            except Exception:
                # A malformed shape shouldn't sink the whole slide.
                continue
        # Speaker notes often carry the real narrative — include them if present.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker notes] {notes}")
        except Exception:
            pass
        text = _clean("\n".join(parts))
        if text:
            sections.append({"label": f"Slide {i}", "text": text})
    return sections, len(slides)


def _extract_pdf(data: bytes) -> tuple[list[dict], int]:
    """Returns (sections-with-text, total page count)."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = reader.pages
    sections: list[dict] = []
    for i, page in enumerate(pages, 1):
        try:
            text = _clean(page.extract_text() or "")
        except Exception:
            text = ""
        if text:
            sections.append({"label": f"Page {i}", "text": text})
    return sections, len(pages)


def extract_deck(filename: str, data: bytes) -> dict:
    """Extract structured text from an uploaded .pptx/.pdf deck.

    Raises ``ValueError`` for an unsupported extension or a deck with no
    extractable text (e.g. a scanned, image-only PDF).
    """
    ext = Path(filename or "").suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"Unsupported file type '{ext or '(none)'}'. Upload a .pptx or .pdf deck.")
    if not data:
        raise ValueError("The uploaded file is empty.")

    kind = "pptx" if ext == ".pptx" else "pdf"
    try:
        sections, total_units = _extract_pptx(data) if kind == "pptx" else _extract_pdf(data)
    except Exception as exc:
        logger.exception("Deck extraction failed for %s", filename)
        raise ValueError(f"Could not read the {kind.upper()} file — it may be corrupt or password-protected.") from exc

    if not sections:
        raise ValueError(
            "No text could be extracted from the deck. If it's a scanned/image-only "
            "PDF, upload a text-based version."
        )

    text = "\n\n".join(f"### {s['label']}\n{s['text']}" for s in sections)
    return {
        "kind": kind,
        "filename": filename,
        "sections": sections,          # only slides/pages that had extractable text
        "total_units": total_units,    # total slides/pages in the file (incl. image-only)
        "text": text,
        "char_count": len(text),
    }
